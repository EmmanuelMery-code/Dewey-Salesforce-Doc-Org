"""PowerPoint export for the dashboard designer.

Renders each dashboard widget (text, table, image, or chart) onto its own
slide. Extracted from ``dashboard_exporter.py`` to keep that module focused
on the simpler PDF/PNG/data export helpers.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Dict, List

import matplotlib.pyplot as plt
from matplotlib.figure import Figure
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


class _DashboardPptxMixin:
    """Adds PowerPoint export to ``DashboardExporter``."""

    def export_to_pptx(self, widgets_data: List[Dict[str, Any]], output_path: Path, title: str = "Salesforce Org Dashboard"):
        """Exporte chaque widget sur une slide PowerPoint séparée."""
        prs = Presentation()
        
        # Slide de titre
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = f"Généré le {pd.Timestamp.now().strftime('%d/%m/%Y')}"

        for widget in widgets_data:
            # Vérifier la visibilité
            widget_data = widget.get('data')
            if isinstance(widget_data, dict) and not widget_data.get('visible', True):
                continue

            # Utiliser un layout avec titre et contenu
            slide_layout = prs.slide_layouts[5] # Title Only
            slide = prs.slides.add_slide(slide_layout)
            
            # Titre de la slide
            slide.shapes.title.text = widget.get('title', 'Composant')
            
            if widget.get('type') == 'text':
                # Ajouter le texte directement dans la slide
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                # Appliquer la couleur de fond
                bg_color = widget.get('color', '#ffffff')
                if bg_color == "none":
                    txBox.fill.background()
                elif bg_color.startswith('#'):
                    from pptx.dml.color import RGBColor
                    r, g, b = int(bg_color[1:3], 16), int(bg_color[3:5], 16), int(bg_color[5:7], 16)
                    txBox.fill.solid()
                    txBox.fill.fore_color.rgb = RGBColor(r, g, b)

                tf = txBox.text_frame
                tf.word_wrap = True
                
                # Alignement vertical global
                from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
                valign = widget.get('text_valign', 'top')
                if valign == 'center': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                elif valign == 'bottom': tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                else: tf.vertical_anchor = MSO_ANCHOR.TOP

                rich = widget.get('rich_text')
                align = widget.get('text_align', 'left')
                pp_align = PP_ALIGN.LEFT
                if align == 'center': pp_align = PP_ALIGN.CENTER
                elif align == 'right': pp_align = PP_ALIGN.RIGHT

                if rich:
                    # On regroupe par paragraphes (séparés par \n dans les segments)
                    tf.clear() # Supprimer le paragraphe par défaut
                    current_p = tf.add_paragraph()
                    current_p.alignment = pp_align
                    
                    for seg in rich:
                        text_parts = seg['text'].split('\n')
                        for i, part in enumerate(text_parts):
                            if i > 0:
                                current_p = tf.add_paragraph()
                                current_p.alignment = pp_align
                            
                            if part:
                                run = current_p.add_run()
                                run.text = part
                                if seg.get('bold'): run.font.bold = True
                                if seg.get('italic'): run.font.italic = True
                                if seg.get('underline'): run.font.underline = True
                                if seg.get('strikeout'): run.font.strike = True
                                if seg.get('size'): run.font.size = Pt(seg['size'])
                                if seg.get('color'):
                                    c = seg['color']
                                    if c.startswith('#'):
                                        from pptx.dml.color import RGBColor
                                        r, g, b = int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)
                                        run.font.color.rgb = RGBColor(r, g, b)
                else:
                    tf.text = widget.get('text', '')
                    tf.paragraphs[0].alignment = pp_align
            elif widget.get('type') == 'table':
                rows_data = widget.get('data', {}).get('rows', [])
                if rows_data:
                    cols_to_show = widget.get('table_columns', [])
                    row_keys = widget.get('table_rows', [])
                    if not cols_to_show:
                        cols_to_show = [k for k in rows_data[0].keys() if k not in row_keys]
                    all_cols = row_keys + cols_to_show
                    
                    # Vérification de la validité des colonnes (renommage via AS)
                    actual_keys = list(rows_data[0].keys())
                    valid_cols = [c for c in all_cols if c in actual_keys]
                    if not valid_cols:
                        all_cols = actual_keys
                    else:
                        all_cols = valid_cols
                    
                    rows_count = min(len(rows_data) + 1, 25) # Limiter pour PPTX
                    cols_count = len(all_cols)
                    
                    table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * rows_count))
                    table = table_shape.table
                    
                    font_name = widget.get('table_font_name', 'Arial')
                    font_size = widget.get('table_font_size', 10)

                    for c, col_name in enumerate(all_cols):
                        cell = table.cell(0, c)
                        cell.text = col_name
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = font_name
                                run.font.size = Pt(font_size + 2) # En-tête un peu plus grand
                                run.font.bold = True
                    
                    for r, row_dict in enumerate(rows_data[:rows_count-1]):
                        for c, col_name in enumerate(all_cols):
                            cell = table.cell(r + 1, c)
                            cell.text = str(row_dict.get(col_name, ''))
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = font_name
                                    run.font.size = Pt(font_size)
            elif widget.get('type') == 'image':
                # Appliquer la couleur de fond
                bg_color = widget.get('color', '#ffffff')
                if bg_color == "none":
                    slide.background.fill.background()
                elif bg_color.startswith('#'):
                    from pptx.dml.color import RGBColor
                    r, g, b = int(bg_color[1:3], 16), int(bg_color[3:5], 16), int(bg_color[5:7], 16)
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

                img_path = widget.get('image_path')
                emoji = widget.get('emoji')
                
                if img_path and Path(img_path).exists():
                    slide.shapes.add_picture(str(img_path), Inches(2), Inches(2), width=Inches(6))
                elif emoji:
                    txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(4))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = emoji
                    p.font.size = Pt(widget.get('image_font_size', 100))
                    
                    font_color = widget.get('image_font_color', '#000000')
                    if font_color.startswith('#'):
                        from pptx.dml.color import RGBColor
                        r, g, b = int(font_color[1:3], 16), int(font_color[3:5], 16), int(font_color[5:7], 16)
                        p.font.color.rgb = RGBColor(r, g, b)
                        
                    from pptx.enum.text import PP_ALIGN
                    p.alignment = PP_ALIGN.CENTER
            else:
                # Générer l'image du graphique pour ce widget seul
                fig = Figure(figsize=(8, 5), dpi=100)
                ax = fig.add_subplot(111)
                self._render_single_widget(ax, widget)
                
                image_stream = io.BytesIO()
                fig.savefig(image_stream, format='png', bbox_inches='tight', dpi=200)
                image_stream.seek(0)
                
                # Centrer l'image
                slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))
                plt.close(fig)

        prs.save(output_path)
        self.log(f"Dashboard exporté en PowerPoint (multi-slides) : {output_path}")
