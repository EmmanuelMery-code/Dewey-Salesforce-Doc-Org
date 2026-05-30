from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Dict, List

import matplotlib.pyplot as plt
from matplotlib.figure import Figure
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
from openpyxl.drawing.image import Image as OpenpyxlImage

class DashboardExporter:
    """Service pour l'exportation des tableaux de bord vers différents formats.
    
    Supporte :
    - PDF (via Matplotlib)
    - PNG (avec ou sans transparence)
    - PowerPoint (PPTX)
    - Excel (XLSX)
    - CSV
    """

    def __init__(self, log_callback=None):
        self.log = log_callback or (lambda msg: None)
        # Configuration globale des polices Matplotlib pour supporter les émojis
        try:
            import matplotlib.pyplot as plt
            emoji_fonts = ["Segoe UI Emoji", "Apple Color Emoji", "Noto Color Emoji", "Symbola", "Segoe UI Symbol"]
            current_sans = list(plt.rcParams.get('font.sans-serif', []))
            for f in reversed(emoji_fonts):
                if f not in current_sans:
                    current_sans.insert(0, f)
            plt.rcParams['font.sans-serif'] = current_sans
            plt.rcParams['font.family'] = 'sans-serif'
            plt.rcParams['mathtext.fontset'] = 'cm'
            plt.rcParams['axes.unicode_minus'] = False
        except: pass

    def export_to_pdf(self, fig: Figure, output_path: Path):
        """Exporte la figure Matplotlib en PDF."""
        fig.savefig(output_path, format='pdf', bbox_inches='tight')
        self.log(f"Dashboard exporté en PDF : {output_path}")

    def export_to_png(self, fig: Figure, output_path: Path, transparent: bool = False):
        """Exporte la figure Matplotlib en PNG."""
        fig.savefig(output_path, format='png', bbox_inches='tight', transparent=transparent)
        self.log(f"Dashboard exporté en PNG ({'transparent' if transparent else 'opaque'}) : {output_path}")

    def export_to_pptx(self, widgets_data: List[Dict[str, Any]], output_path: Path, title: str = "Salesforce Org Dashboard"):
        """Exporte chaque widget sur une slide PowerPoint séparée."""
        prs = Presentation()
        
        # Slide de titre
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = f"Généré le {pd.Timestamp.now().strftime('%d/%m/%Y')}"

        for widget in widgets_data:
            # Vérifier la visibilité
            widget_data = widget.get('data')
            if isinstance(widget_data, dict) and not widget_data.get('visible', True):
                continue

            # Utiliser un layout avec titre et contenu
            slide_layout = prs.slide_layouts[5] # Title Only
            slide = prs.slides.add_slide(slide_layout)
            
            # Titre de la slide
            slide.shapes.title.text = widget.get('title', 'Composant')
            
            if widget.get('type') == 'text':
                # Ajouter le texte directement dans la slide
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                # Appliquer la couleur de fond
                bg_color = widget.get('color', '#ffffff')
                if bg_color == "none":
                    txBox.fill.background()
                elif bg_color.startswith('#'):
                    from pptx.dml.color import RGBColor
                    r, g, b = int(bg_color[1:3], 16), int(bg_color[3:5], 16), int(bg_color[5:7], 16)
                    txBox.fill.solid()
                    txBox.fill.fore_color.rgb = RGBColor(r, g, b)

                tf = txBox.text_frame
                tf.word_wrap = True
                
                # Alignement vertical global
                from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
                valign = widget.get('text_valign', 'top')
                if valign == 'center': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                elif valign == 'bottom': tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                else: tf.vertical_anchor = MSO_ANCHOR.TOP

                rich = widget.get('rich_text')
                align = widget.get('text_align', 'left')
                pp_align = PP_ALIGN.LEFT
                if align == 'center': pp_align = PP_ALIGN.CENTER
                elif align == 'right': pp_align = PP_ALIGN.RIGHT

                if rich:
                    # On regroupe par paragraphes (séparés par \n dans les segments)
                    tf.clear() # Supprimer le paragraphe par défaut
                    current_p = tf.add_paragraph()
                    current_p.alignment = pp_align
                    
                    for seg in rich:
                        text_parts = seg['text'].split('\n')
                        for i, part in enumerate(text_parts):
                            if i > 0:
                                current_p = tf.add_paragraph()
                                current_p.alignment = pp_align
                            
                            if part:
                                run = current_p.add_run()
                                run.text = part
                                if seg.get('bold'): run.font.bold = True
                                if seg.get('italic'): run.font.italic = True
                                if seg.get('underline'): run.font.underline = True
                                if seg.get('strikeout'): run.font.strike = True
                                if seg.get('size'): run.font.size = Pt(seg['size'])
                                if seg.get('color'):
                                    c = seg['color']
                                    if c.startswith('#'):
                                        from pptx.dml.color import RGBColor
                                        r, g, b = int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)
                                        run.font.color.rgb = RGBColor(r, g, b)
                else:
                    tf.text = widget.get('text', '')
                    tf.paragraphs[0].alignment = pp_align
            elif widget.get('type') == 'table':
                rows_data = widget.get('data', {}).get('rows', [])
                if rows_data:
                    cols_to_show = widget.get('table_columns', [])
                    row_keys = widget.get('table_rows', [])
                    if not cols_to_show:
                        cols_to_show = [k for k in rows_data[0].keys() if k not in row_keys]
                    all_cols = row_keys + cols_to_show
                    
                    # Vérification de la validité des colonnes (renommage via AS)
                    actual_keys = list(rows_data[0].keys())
                    valid_cols = [c for c in all_cols if c in actual_keys]
                    if not valid_cols:
                        all_cols = actual_keys
                    else:
                        all_cols = valid_cols
                    
                    rows_count = min(len(rows_data) + 1, 25) # Limiter pour PPTX
                    cols_count = len(all_cols)
                    
                    table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * rows_count))
                    table = table_shape.table
                    
                    font_name = widget.get('table_font_name', 'Arial')
                    font_size = widget.get('table_font_size', 10)

                    for c, col_name in enumerate(all_cols):
                        cell = table.cell(0, c)
                        cell.text = col_name
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = font_name
                                run.font.size = Pt(font_size + 2) # En-tête un peu plus grand
                                run.font.bold = True
                    
                    for r, row_dict in enumerate(rows_data[:rows_count-1]):
                        for c, col_name in enumerate(all_cols):
                            cell = table.cell(r + 1, c)
                            cell.text = str(row_dict.get(col_name, ''))
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = font_name
                                    run.font.size = Pt(font_size)
            elif widget.get('type') == 'image':
                # Appliquer la couleur de fond
                bg_color = widget.get('color', '#ffffff')
                if bg_color == "none":
                    slide.background.fill.background()
                elif bg_color.startswith('#'):
                    from pptx.dml.color import RGBColor
                    r, g, b = int(bg_color[1:3], 16), int(bg_color[3:5], 16), int(bg_color[5:7], 16)
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

                img_path = widget.get('image_path')
                emoji = widget.get('emoji')
                
                if img_path and Path(img_path).exists():
                    slide.shapes.add_picture(str(img_path), Inches(2), Inches(2), width=Inches(6))
                elif emoji:
                    txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(4))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = emoji
                    p.font.size = Pt(widget.get('image_font_size', 100))
                    
                    font_color = widget.get('image_font_color', '#000000')
                    if font_color.startswith('#'):
                        from pptx.dml.color import RGBColor
                        r, g, b = int(font_color[1:3], 16), int(font_color[3:5], 16), int(font_color[5:7], 16)
                        p.font.color.rgb = RGBColor(r, g, b)
                        
                    from pptx.enum.text import PP_ALIGN
                    p.alignment = PP_ALIGN.CENTER
            else:
                # Générer l'image du graphique pour ce widget seul
                fig = Figure(figsize=(8, 5), dpi=100)
                ax = fig.add_subplot(111)
                self._render_single_widget(ax, widget)
                
                image_stream = io.BytesIO()
                fig.savefig(image_stream, format='png', bbox_inches='tight', dpi=200)
                image_stream.seek(0)
                
                # Centrer l'image
                slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))
                plt.close(fig)

        prs.save(output_path)
        self.log(f"Dashboard exporté en PowerPoint (multi-slides) : {output_path}")

    def _render_single_widget(self, ax, widget):
        """Helper pour rendre un seul widget sur un axe donné."""
        data = widget.get('data', {})
        w_type = widget.get('type')
        color = widget.get('color', '#3498db')
        
        # Gestion intelligente de la couleur de fond
        if w_type in ["text", "image", "kpi"] and ',' not in color:
            bg_color = color
        else:
            bg_color = "none"

        if bg_color == "none":
            ax.set_facecolor('none')
            if hasattr(ax, 'figure'):
                ax.figure.patch.set_alpha(0.0)
        else:
            try:
                ax.set_facecolor(bg_color)
            except:
                ax.set_facecolor('none')

        if w_type == "text":
            ax.axis('off')
            if color != "none":
                from matplotlib.patches import Rectangle
                rect = Rectangle((0,0), 1, 1, transform=ax.transAxes, color=color, zorder=-1)
                ax.add_patch(rect)
            
            rich = widget.get('rich_text')
            align = widget.get('text_align', 'left')
            valign = widget.get('text_valign', 'top')
            
            if rich:
                from matplotlib.offsetbox import TextArea, HPacker, AnnotationBbox, VPacker
                lines = []
                current_line_segments = []
                
                for seg in rich:
                    text_parts = seg['text'].split('\n')
                    for i, part in enumerate(text_parts):
                        if i > 0:
                            lines.append(current_line_segments)
                            current_line_segments = []
                        if part:
                            props = {
                                "color": seg.get("color", "black"),
                                "fontsize": seg.get("size", 10),
                                "fontweight": "bold" if seg.get("bold") else "normal",
                                "fontstyle": "italic" if seg.get("italic") else "normal",
                                "fontname": seg.get("font", "Arial")
                            }
                            
                            # Support des émojis dans le texte riche
                            if any(ord(c) > 0xFFFF for c in part):
                                # On laisse le fallback global gérer l'émoji
                                props["fontname"] = "sans-serif"
                                props["fontweight"] = "normal"
                                props["fontstyle"] = "normal"

                            current_line_segments.append(TextArea(part, textprops=props))
                        elif i > 0:
                            # Ligne vide : on ajoute un espace pour maintenir la hauteur
                            current_line_segments.append(TextArea(" "))
                if current_line_segments: lines.append(current_line_segments)
                if lines:
                    line_boxes = [HPacker(children=l, align="baseline", pad=0, sep=0) for l in lines]
                    vbox = VPacker(children=line_boxes, align=align, pad=0, sep=2)
                    
                    xy = (0.05, 0.95)
                    box_align = (0, 1)
                    if align == "center": xy = (0.5, xy[1]); box_align = (0.5, box_align[1])
                    elif align == "right": xy = (0.95, xy[1]); box_align = (1, box_align[1])
                    
                    if valign == "center": xy = (xy[0], 0.5); box_align = (box_align[0], 0.5)
                    elif valign == "bottom": xy = (xy[0], 0.05); box_align = (box_align[0], 0)
                    
                    ab = AnnotationBbox(vbox, xy, xycoords='axes fraction', box_alignment=box_align, frameon=False)
                    ax.add_artist(ab)
            else:
                ha = align
                va = valign
                xy = (0.05, 0.95)
                if ha == "center": xy = (0.5, 0.5 if va == "center" else (0.05 if va == "bottom" else 0.95))
                elif ha == "right": xy = (0.95, 0.5 if va == "center" else (0.05 if va == "bottom" else 0.95))
                # Support des émojis pour le texte simple
                ax.text(xy[0], xy[1], widget.get('text', ''), va=va, ha=ha, wrap=True, fontsize=10, transform=ax.transAxes)
            return

        # Nettoyer les données pour le rendu (enlever la clé technique 'visible')
        plot_data = {k: v for k, v in data.items() if k != 'visible'} if isinstance(data, dict) else data
        
        # LOGGING pour débogage export
        print(f"Export widget ({w_type}): data={plot_data}")

        if not plot_data and w_type != "image":
            ax.axis('off')
            ax.text(0.5, 0.5, "Aucune donnée", ha='center', va='center')
            return

        if w_type == "pie" or w_type == "donut":
            ax.axis('off')
            valid_data = {str(k): float(v) for k, v in plot_data.items() 
                         if isinstance(v, (int, float, str)) and str(v).replace('.','',1).isdigit() and float(v) > 0}
            
            labels = list(valid_data.keys())
            values = list(valid_data.values())
            
            if not values:
                ax.text(0.5, 0.5, "Données nulles", ha='center', va='center')
            else:
                colors = color.split(',') if ',' in color else None
                if w_type == "pie":
                    ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=140, colors=colors)
                else:
                    ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=140, colors=colors, wedgeprops=dict(width=0.4))
        elif w_type == "bar":
            labels = [str(k) for k in plot_data.keys()]
            values = [float(v) if isinstance(v, (int, float, str)) and str(v).replace('.','',1).isdigit() else 0.0 for v in plot_data.values()]
            if not values or all(v == 0 for v in values):
                ax.text(0.5, 0.5, "Données nulles", ha='center', va='center')
            else:
                try:
                    colors = color.split(',') if ',' in color else color
                    x_indices = list(range(len(labels)))
                    ax.bar(x_indices, values, color=colors)
                    ax.set_xticks(x_indices)
                    ax.set_xticklabels(labels, rotation=45, ha='right', fontsize=8)
                except Exception as e:
                    print(f"Erreur export bar: {e}")
                    ax.text(0.5, 0.5, "Erreur graphique", ha='center', va='center')
        elif w_type == "stacked_bar":
            labels = plot_data.get('labels', [])
            series = plot_data.get('series', {})
            if not labels or not series:
                ax.text(0.5, 0.5, "Données nulles", ha='center', va='center')
            else:
                try:
                    x_indices = list(range(len(labels)))
                    bottom = None
                    colors = color.split(',') if ',' in color else [color]
                    for i, (name, vals) in enumerate(series.items()):
                        c = colors[i % len(colors)]
                        clean_vals = [float(v) if isinstance(v, (int, float, str)) and str(v).replace('.','',1).isdigit() else 0.0 for v in vals]
                        ax.bar(x_indices, clean_vals, bottom=bottom, label=name, color=c)
                        if bottom is None: bottom = [0.0] * len(clean_vals)
                        bottom = [b + v for b, v in zip(bottom, clean_vals)]
                    
                    ax.legend(fontsize=7)
                    ax.set_xticks(x_indices)
                    step = max(1, len(labels) // 6)
                    display_labels = [l if i % step == 0 else "" for i, l in enumerate(labels)]
                    ax.set_xticklabels(display_labels, rotation=45, ha='right', fontsize=8)
                except Exception as e:
                    print(f"Erreur export stacked_bar: {e}")
                    ax.text(0.5, 0.5, "Erreur graphique", ha='center', va='center')
        elif w_type == "line" or w_type == "area":
            labels = [str(k) for k in plot_data.keys()]
            def safe_float(v):
                try:
                    if isinstance(v, str): v = v.replace(',', '.')
                    return float(v)
                except: return 0.0
            values = [safe_float(v) for v in plot_data.values()]
            if not labels:
                ax.text(0.5, 0.5, "Pas de données", ha='center', va='center')
            else:
                try:
                    line_color = color.split(',')[0] if ',' in color else color
                    if line_color == "none": line_color = "#3498db"
                    x_indices = list(range(len(labels)))
                    
                    if w_type == "line":
                        ax.plot(x_indices, values, marker='o', color=line_color, linewidth=2)
                    else:
                        ax.fill_between(x_indices, values, color=line_color, alpha=0.3)
                        ax.plot(x_indices, values, color=line_color, marker='.', linewidth=1)
                    
                    ax.set_xticks(x_indices)
                    step = max(1, len(labels) // 5)
                    display_labels = [l if i % step == 0 else "" for i, l in enumerate(labels)]
                    ax.set_xticklabels(display_labels, rotation=30, ha='right', fontsize=8)
                    ax.grid(True, linestyle='--', alpha=0.6)
                except Exception as e:
                    print(f"Erreur export line/area: {e}")
                    ax.text(0.5, 0.5, "Erreur graphique", ha='center', va='center')
        elif w_type == "kpi":
            ax.axis('off')
            # Formater les nombres
            dec = widget.get('kpi_decimals', 1)
            lines = []
            for k, v in plot_data.items():
                if isinstance(v, (int, float)):
                    lines.append(f"{k}: {v:.{dec}f}")
                else:
                    lines.append(f"{k}: {v}")
            text = "\n".join(lines)
            
            bbox_props = dict(facecolor='white', alpha=0.3, boxstyle='round')
            if color == "none":
                bbox_props['alpha'] = 0.0
                
            ax.text(0.5, 0.5, text, ha='center', va='center', fontsize=16, fontweight='bold', 
                    bbox=bbox_props)
        elif w_type == "table":
            ax.axis('off')
            rows = data.get('rows', [])
            if rows:
                cols_to_show = widget.get('table_columns', [])
                row_keys = widget.get('table_rows', [])
                if not cols_to_show:
                    cols_to_show = [k for k in rows[0].keys() if k not in row_keys]
                all_cols = row_keys + cols_to_show
                
                # Vérification de la validité des colonnes (renommage via AS)
                actual_keys = list(rows[0].keys())
                valid_cols = [c for c in all_cols if c in actual_keys]
                if not valid_cols:
                    all_cols = actual_keys
                else:
                    all_cols = valid_cols

                table_data = [[str(r.get(c, '')) for c in all_cols] for r in rows[:20]]
                tab = ax.table(cellText=table_data, colLabels=all_cols, loc='center', cellLoc='center', bbox=[0, 0, 1, 1])
                tab.auto_set_font_size(False)
                
                font_name = widget.get('table_font_name', 'Segoe UI')
                font_size = widget.get('table_font_size', 8)
                tab.set_fontsize(font_size)
                
                for cell in tab.get_celld().values():
                    cell.set_text_props(fontfamily=font_name)
                
                tab.auto_set_column_width(col=list(range(len(all_cols))))
        elif w_type == "image":
            ax.axis('off')
            if color != "none":
                from matplotlib.patches import Rectangle
                rect = Rectangle((0,0), 1, 1, transform=ax.transAxes, color=color, zorder=-1)
                ax.add_patch(rect)
            
            img_path = widget.get('image_path')
            emoji = widget.get('emoji')
            
            if img_path and Path(img_path).exists():
                try:
                    import matplotlib.image as mpimg
                    img = mpimg.imread(img_path)
                    ax.imshow(img, aspect='equal', extent=[0.1, 0.9, 0.1, 0.9])
                except:
                    ax.text(0.5, 0.5, "Erreur image", ha='center', va='center')
            elif emoji:
                # On utilise le fallback global via sans-serif
                ax.text(0.5, 0.5, emoji, ha='center', va='center', 
                        fontsize=widget.get('image_font_size', 60), 
                        color=widget.get('image_font_color', '#000000'),
                        fontfamily="sans-serif")
            else:
                ax.text(0.5, 0.5, "Aucune image", ha='center', va='center')

    def export_data(self, widgets_data: List[Dict[str, Any]], output_path: Path, format: str = 'excel'):
        """Exporte les données brutes des widgets en Excel ou CSV."""
        # On aplatit les données pour l'export
        flat_data = []
        for widget in widgets_data:
            w_name = widget.get('title', 'Sans titre')
            w_data = widget.get('data', {})
            
            if isinstance(w_data, dict) and "rows" in w_data:
                for row_dict in w_data["rows"]:
                    row = {'Widget': w_name}
                    row.update(row_dict)
                    flat_data.append(row)
            elif isinstance(w_data, dict):
                for k, v in w_data.items():
                    flat_data.append({'Widget': w_name, 'Metrique': k, 'Valeur': v})
            elif isinstance(w_data, list):
                for item in w_data:
                    if isinstance(item, dict):
                        row = {'Widget': w_name}
                        row.update(item)
                        flat_data.append(row)
                    else:
                        flat_data.append({'Widget': w_name, 'Valeur': item})

        df = pd.DataFrame(flat_data)
        
        if format == 'excel':
            df.to_excel(output_path, index=False)
            self.log(f"Données du dashboard exportées en Excel : {output_path}")
        else:
            df.to_csv(output_path, index=False, encoding='utf-8-sig')
            self.log(f"Données du dashboard exportées en CSV : {output_path}")
